from datetime import datetime
import os
import uuid
import tiktoken
from typing import List, Dict, Any, Optional
from pathlib import Path
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
from openai import OpenAI
from sqlalchemy.orm import Session
from sqlalchemy import text

from app.core.config import settings
from app.models.brand_document import BrandDocument, DocumentChunk

# Directory per upload
UPLOAD_DIR = Path("/var/www/noscite-calendar/uploads/documents")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

# Tokenizer per contare token
encoding = tiktoken.get_encoding("cl100k_base")


class RAGService:
    def __init__(self):
        self.openai_client = OpenAI(api_key=settings.OPENAI_API_KEY)
        self.chunk_size = 500
        self.chunk_overlap = 50
    
    # === FILE EXTRACTION ===
    
    def extract_text_from_pdf(self, file_path: str) -> str:
        text = ""
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n\n"
        except Exception as e:
            print(f"Errore estrazione PDF: {e}")
        return text.strip()
    
    def extract_text_from_docx(self, file_path: str) -> str:
        text = ""
        try:
            doc = DocxDocument(file_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n\n"
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        text += row_text + "\n"
                text += "\n"
        except Exception as e:
            print(f"Errore estrazione DOCX: {e}")
        return text.strip()
    
    def extract_text_from_pptx(self, file_path: str) -> str:
        text = ""
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"--- Slide {slide_num} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                if slide_text.strip() != f"--- Slide {slide_num} ---":
                    text += slide_text + "\n"
        except Exception as e:
            print(f"Errore estrazione PPTX: {e}")
        return text.strip()
    
    def extract_text_from_txt(self, file_path: str) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
    
    def extract_text(self, file_path: str, file_type: str) -> str:
        extractors = {
            'pdf': self.extract_text_from_pdf,
            'docx': self.extract_text_from_docx,
            'doc': self.extract_text_from_docx,
            'pptx': self.extract_text_from_pptx,
            'ppt': self.extract_text_from_pptx,
            'txt': self.extract_text_from_txt,
            'md': self.extract_text_from_txt,
        }
        extractor = extractors.get(file_type.lower())
        if extractor:
            return extractor(file_path)
        return ""
    
    # === CHUNKING ===
    
    def count_tokens(self, text: str) -> int:
        return len(encoding.encode(text))
    
    def chunk_text(self, text: str) -> List[Dict[str, Any]]:
        chunks = []
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        current_chunk = ""
        current_tokens = 0
        chunk_index = 0
        
        for para in paragraphs:
            para_tokens = self.count_tokens(para)
            
            if para_tokens > self.chunk_size:
                if current_chunk:
                    chunks.append({
                        "index": chunk_index,
                        "content": current_chunk.strip(),
                        "token_count": current_tokens
                    })
                    chunk_index += 1
                    current_chunk = ""
                    current_tokens = 0
                
                sentences = para.replace('. ', '.|').split('|')
                for sentence in sentences:
                    sent_tokens = self.count_tokens(sentence)
                    if current_tokens + sent_tokens > self.chunk_size:
                        if current_chunk:
                            chunks.append({
                                "index": chunk_index,
                                "content": current_chunk.strip(),
                                "token_count": current_tokens
                            })
                            chunk_index += 1
                        current_chunk = sentence
                        current_tokens = sent_tokens
                    else:
                        current_chunk += " " + sentence
                        current_tokens += sent_tokens
            
            elif current_tokens + para_tokens > self.chunk_size:
                if current_chunk:
                    chunks.append({
                        "index": chunk_index,
                        "content": current_chunk.strip(),
                        "token_count": current_tokens
                    })
                    chunk_index += 1
                
                overlap_text = current_chunk[-200:] if len(current_chunk) > 200 else ""
                current_chunk = overlap_text + "\n\n" + para
                current_tokens = self.count_tokens(current_chunk)
            else:
                current_chunk += "\n\n" + para
                current_tokens += para_tokens
        
        if current_chunk.strip():
            chunks.append({
                "index": chunk_index,
                "content": current_chunk.strip(),
                "token_count": self.count_tokens(current_chunk)
            })
        
        return chunks
    
    # === EMBEDDINGS ===
    
    def generate_embedding(self, text: str) -> List[float]:
        response = self.openai_client.embeddings.create(
            model="text-embedding-3-small",
            input=text
        )
        return response.data[0].embedding
    
    def generate_embeddings_batch(self, texts: List[str]) -> List[List[float]]:
        response = self.openai_client.embeddings.create(
            model="text-embedding-3-small",
            input=texts
        )
        return [item.embedding for item in response.data]
    
    # === DOCUMENT ANALYSIS ===
    
    async def analyze_document(self, document: BrandDocument, text: str, db: Session) -> Dict[str, Any]:
        """Analizza il documento con AI per estrarre summary, tipo e key topics"""
        from anthropic import Anthropic
        
        client = Anthropic(api_key=settings.ANTHROPIC_API_KEY)
        text_sample = text[:8000] if len(text) > 8000 else text
        
        prompt = f"""Analizza questo documento aziendale e fornisci:

1. **Tipo documento**: Classifica in UNA categoria:
   - brand_guidelines, company_presentation, product_info, case_study
   - marketing_material, internal_docs, blog_content, press_release, other

2. **Riassunto**: 2-3 frasi del contenuto principale.

3. **Key Topics**: 3-5 argomenti/parole chiave principali.

4. **Tone of Voice**: Descrivi il tono (formale, informale, tecnico, ecc.)

5. **Target Audience**: A chi è rivolto.

DOCUMENTO:
{text_sample}

Rispondi SOLO in JSON:
{{"document_type": "tipo", "summary": "riassunto", "key_topics": ["t1", "t2"], "tone_of_voice": "tono", "target_audience": "target"}}"""

        try:
            response = client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=1000,
                messages=[{"role": "user", "content": prompt}]
            )
            
            response_text = response.content[0].text
            if "```json" in response_text:
                response_text = response_text.split("```json")[1].split("```")[0]
            elif "```" in response_text:
                response_text = response_text.split("```")[1].split("```")[0]
            
            import json
            analysis = json.loads(response_text.strip())
            
            document.summary = analysis.get("summary", "")
            document.key_topics = {
                "topics": analysis.get("key_topics", []),
                "document_type": analysis.get("document_type", "other"),
                "tone_of_voice": analysis.get("tone_of_voice", ""),
                "target_audience": analysis.get("target_audience", "")
            }
            document.analysis_status = "completed"
            document.analyzed_at = datetime.now()
            db.commit()
            
            return analysis
        except Exception as e:
            print(f"Errore analisi documento: {e}")
            document.analysis_status = "failed"
            db.commit()
            return {}
    
    # === DOCUMENT PROCESSING ===
    
    async def process_document(self, document: BrandDocument, db: Session) -> bool:
        """Processa un documento: estrae testo, chunka, genera embeddings, analizza con AI"""
        try:
            document.extraction_status = "processing"
            db.commit()
            
            text = self.extract_text(document.file_path, document.file_type)
            
            if not text:
                document.extraction_status = "failed"
                db.commit()
                return False
            
            document.extraction_status = "completed"
            db.commit()
            
            # Analisi AI
            document.analysis_status = "processing"
            db.commit()
            await self.analyze_document(document, text, db)
            
            # Chunking
            chunks = self.chunk_text(text)
            if not chunks:
                return True
            
            # Embeddings in batch
            chunk_texts = [c["content"] for c in chunks]
            all_embeddings = []
            for i in range(0, len(chunk_texts), 100):
                batch = chunk_texts[i:i+100]
                embeddings = self.generate_embeddings_batch(batch)
                all_embeddings.extend(embeddings)
            
            # Salva chunks
            for chunk, embedding in zip(chunks, all_embeddings):
                db_chunk = DocumentChunk(
                    document_id=document.id,
                    brand_id=document.brand_id,
                    chunk_index=chunk["index"],
                    content=chunk["content"],
                    token_count=chunk["token_count"],
                    embedding=embedding
                )
                db.add(db_chunk)
            
            db.commit()
            return True
            
        except Exception as e:
            print(f"Errore processing documento: {e}")
            document.extraction_status = "failed"
            document.analysis_status = "failed"
            db.commit()
            return False
    
    # === SEMANTIC SEARCH ===
    
    def search_similar(self, db: Session, brand_id: int, query: str, limit: int = 5) -> List[Dict[str, Any]]:
        query_embedding = self.generate_embedding(query)
        
        sql = text("""
            SELECT 
                dc.id, dc.content, dc.chunk_index, dc.document_id,
                bd.original_filename,
                1 - (dc.embedding <=> :embedding::vector) as similarity
            FROM document_chunks dc
            JOIN brand_documents bd ON dc.document_id = bd.id
            WHERE dc.brand_id = :brand_id
            ORDER BY dc.embedding <=> :embedding::vector
            LIMIT :limit
        """)
        
        result = db.execute(sql, {
            "brand_id": brand_id,
            "embedding": str(query_embedding),
            "limit": limit
        })
        
        return [
            {
                "id": row[0],
                "content": row[1],
                "chunk_index": row[2],
                "document_id": row[3],
                "filename": row[4],
                "similarity": float(row[5])
            }
            for row in result
        ]
    
    def get_context_for_generation(self, db: Session, brand_id: int, topic: str, max_tokens: int = 2000) -> str:
        chunks = self.search_similar(db, brand_id, topic, limit=10)
        
        context_parts = []
        total_tokens = 0
        
        for chunk in chunks:
            chunk_tokens = self.count_tokens(chunk["content"])
            if total_tokens + chunk_tokens > max_tokens:
                break
            context_parts.append(f"[Da: {chunk['filename']}]\n{chunk['content']}")
            total_tokens += chunk_tokens
        
        return "\n\n---\n\n".join(context_parts)


# Singleton
rag_service = RAGService()
